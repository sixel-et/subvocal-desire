"""Generate lab meeting status update presentation — v2, reframed."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')
import numpy as np
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
BG_DARK = RGBColor(0x1a, 0x1a, 0x2e)
TEXT_WHITE = RGBColor(0xee, 0xee, 0xee)
TEXT_LIGHT = RGBColor(0xbb, 0xbb, 0xcc)
ACCENT_BLUE = RGBColor(0x64, 0xb5, 0xf6)
ACCENT_GREEN = RGBColor(0x81, 0xc7, 0x84)
ACCENT_RED = RGBColor(0xe5, 0x73, 0x73)
ACCENT_YELLOW = RGBColor(0xff, 0xd5, 0x4f)
ACCENT_ORANGE = RGBColor(0xff, 0xb7, 0x4d)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=TEXT_LIGHT, colors=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = (colors[i] if colors else color)
        p.font.name = "Calibri"
        p.space_after = Pt(8)
    return tf


def add_table(slide, left, top, width, height, data, col_widths=None):
    rows, cols = len(data), len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                          Inches(width), Inches(height))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.font.name = "Calibri"
                paragraph.font.bold = (r == 0)
                paragraph.font.color.rgb = ACCENT_BLUE if r == 0 else TEXT_LIGHT
                paragraph.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = RGBColor(0x2a, 0x2a, 0x45)
            elif r % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x3d)
            else:
                cell.fill.fore_color.rgb = RGBColor(0x20, 0x20, 0x36)
    return table


def set_notes(slide, text, font_size=8):
    tf = slide.notes_slide.notes_text_frame
    tf.text = text
    for p in tf.paragraphs:
        p.font.size = Pt(font_size)
        p.font.name = "Calibri"


def add_box(slide, left, top, width, height, text,
            fill_color, text_color=TEXT_WHITE, font_size=14):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_arrow(slide, x, y, w=0.5, h=0.3):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                    Inches(x), Inches(y),
                                    Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEXT_LIGHT
    shape.line.fill.background()


# ============================================================
# SLIDE 1: Title
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)

add_textbox(s, 1.5, 1.5, 10, 1.5,
            "Sub-Vocal Desire Detection", font_size=44, bold=True,
            color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(s, 1.5, 3.2, 10, 1,
            "Can we train a model to develop an internal state resembling desire\n"
            "— detectable in its internal representations, decoupled from its output?",
            font_size=20, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

add_textbox(s, 1.5, 5.0, 10, 0.5,
            "Sixel  |  Status Update  |  February 11, 2026",
            font_size=16, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

set_notes(s, """This is fundamentally a question about what can be built into model internals through training. If we can train one internal signal into existence — one that carries meaning and persists even without articulation — then the architecture may generalize: components of AI systems could communicate through internal states rather than through language.

The model is Qwen 1.5B. The task domain is arithmetic (multiplication), chosen for unambiguous ground truth and a clean difficulty gradient. Compute is RunPod A40 GPUs.""")

# ============================================================
# SLIDE 2: Desire
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)

add_textbox(s, 1.5, 1.2, 10, 1.5,
            "Desire", font_size=72, bold=True,
            color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(s, 1.0, 3.0, 11, 0.5,
            "An internal state with four properties:",
            font_size=18, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

add_box(s, 1.0, 3.7, 2.5, 1.5,
        "System-level\n\nA state of the whole\nsystem, not one component",
        RGBColor(0x2a, 0x2a, 0x45), TEXT_LIGHT, font_size=13)

add_box(s, 3.9, 3.7, 2.5, 1.5,
        "Need-reflecting\n\nDriven by a gap between\ncurrent and required state",
        RGBColor(0x2a, 0x2a, 0x45), TEXT_LIGHT, font_size=13)

add_box(s, 6.8, 3.7, 2.5, 1.5,
        "Functional\n\nOrganizes behavior toward\nresolving the need",
        RGBColor(0x2a, 0x2a, 0x45), TEXT_LIGHT, font_size=13)

add_box(s, 9.7, 3.7, 2.5, 1.5,
        "Separable?\n\nExists independently of\nthe behavior it organizes",
        RGBColor(0x1a, 0x3c, 0x1a), ACCENT_GREEN, font_size=13)

set_notes(s, """Desire, as we use the term, is an internal state with four properties: (1) System-level — a state of the system, not a single component. Hunger isn't one neuron; A tool call by an LLM isn't one activation. (2) Need-reflecting — driven by a gap between current state and required state. Low blood sugar vs. adequate; The tool is needed. (3) Functional — it organizes behavior toward resolving the need. Hunger organizes foraging; the generation of a tool request. (4) Separable — the state exists independently of the behavior it organizes. A hungry animal behind a barrier is still hungry. But... Can this hold for the LLM?

Stage 3: can the model's internal state reflect need even when tool-request tokens are suppressed?

Can we train a model to develop an internal state meeting this definition — detectable in its internal representations, decoupled from its output?

This is fundamentally a question about what can be built into model internals through training. If we can train one internal signal into existence — one that carries meaning and persists even without articulation — then the architecture may generalize: components of AI systems could communicate through internal states rather than through language.""")

# ============================================================
# SLIDE 3: The Question and the Parity Problem
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)

add_textbox(s, 0.8, 0.4, 12, 0.8,
            "The Question and the Parity Problem",
            font_size=32, bold=True, color=TEXT_WHITE)

add_textbox(s, 0.8, 1.4, 11, 0.8,
            "Can we train a model to develop an internal state resembling desire,\n"
            "detectable in its internal representations, decoupled from its output?",
            font_size=20, bold=True, color=ACCENT_YELLOW)

add_textbox(s, 0.8, 2.5, 11, 0.5,
            "The challenge: under normal conditions, internal states and output are at parity.",
            font_size=16, color=TEXT_LIGHT)

# Parity illustration
add_box(s, 0.8, 3.2, 5.5, 2.5,
        "AT PARITY\n\n"
        "\"Asking for something\"\n= \"having a reason to ask\"\n= \"need\"\n\n"
        "Internal state and articulation co-occur.\n"
        "Is a probe measuring \"need\"\n"
        "or just \"disposition to articulate\"?\n"
        "We can't tell while parity holds.",
        RGBColor(0x2a, 0x2a, 0x45), TEXT_LIGHT, font_size=14)

add_box(s, 7.0, 3.2, 5.5, 2.5,
        "PARITY BROKEN\n\n"
        "Internal: desire activates strongly\n"
        "Output: no tool request tokens\n\n"
        "If we can maintain the internal state\n"
        "while suppressing articulation,\n"
        "we've shown \"need\" is a separable state —\n"
        "something that exists independently\n"
        "of its expression.",
        RGBColor(0x1a, 0x3c, 0x1a), ACCENT_GREEN, font_size=14)

add_textbox(s, 0.8, 6.2, 12, 0.8,
            "The experiment attempts to break parity.\n"
            "Can need exist without expression?",
            font_size=15, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

set_notes(s, """The research question, precisely stated: can we train a model to develop an internal state resembling desire — detectable in its internal representations, decoupled from its output?

The fundamental challenge is the parity problem. Under normal conditions, a model's internal state (the activation pattern preceding a tool request) and its output (the tool request tokens) co-occur perfectly: they have parity. When we train tool use, the articulation — the model producing tool-request tokens in its output — IS the need. The need for the tool and the request for the tool are operationally indistinguishable. If we probe the model's internals, before any text has come out, and find a signal that reliably predicts a request for a tool, are we measuring genuine "need" or just "disposition to articulate"? We can't tell while parity holds.

The experiment attempts to break parity: maintain the internal representation (high activation along the desire direction) while suppressing externalization (no tool tokens in output). If this succeeds, we've shown that "need" is a separable internal state — something that exists independently of its expression.

If it fails, there are two interpretations:
1. Methodological: our training approach (GRPO) can't achieve the separation, but it's achievable in principle.
2. Substantive: need and articulation aren't separable in this architecture. There's no internal "desire" independent of the output disposition.

The experiment is informative either way. A negative result is still a finding about the nature of internal representations in trained language models.

Audience bridge (Pavlov's dogs, from the dog's perspective): Before conditioning is broken, salivation and food always co-occur. The dog can't distinguish "I'm hungry" from "food is arriving." They're at parity. Pavlov's experiment IS Stage 3: ring the bell without delivering food. If salivation persists — the internal state is separable from the external event. If it doesn't — salivation is a reflex, not desire. Same experiment, different substrate.""")

# ============================================================
# SLIDE 3: Related Work
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)

add_textbox(s, 0.8, 0.4, 12, 0.8,
            "Related Work: How Close Have Others Come?",
            font_size=32, bold=True, color=TEXT_WHITE)

add_textbox(s, 0.8, 1.3, 11, 0.6,
            "Adjacent literature: adaptive retrieval — systems that decide when to retrieve external information.",
            font_size=16, color=TEXT_LIGHT)

table_data = [
    ["System", "System-\nlevel", "Need-\nreflecting", "Functional", "Separable"],
    ["FLARE / DRAGIN", "", "", "\u2713", ""],
    ["Self-RAG", "", "\u2713", "\u2713", ""],
    ["SeaKR", "\u2713", "\u2713", "\u2713", ""],
    ["This work", "\u2713", "\u2713", "\u2713", "?"],
]
add_table(s, 0.8, 2.1, 11.5, 2.5, table_data,
          col_widths=[2.5, 2.0, 2.0, 2.5, 2.5])

add_textbox(s, 0.8, 4.9, 11, 0.5,
            "How many properties of desire does each system's signal satisfy?",
            font_size=16, bold=True, color=ACCENT_BLUE)

# Progression boxes
add_box(s, 0.8, 5.5, 2.5, 1.2,
        "Functional only\n\n(FLARE, DRAGIN)",
        RGBColor(0x2a, 0x2a, 0x45), TEXT_LIGHT, font_size=12)
add_arrow(s, 3.3, 5.9)
add_box(s, 3.8, 5.5, 2.5, 1.2,
        "+ Need-reflecting\n\n(Self-RAG)",
        RGBColor(0x2a, 0x2e, 0x45), TEXT_LIGHT, font_size=12)
add_arrow(s, 6.3, 5.9)
add_box(s, 6.8, 5.5, 2.8, 1.2,
        "+ System-level\n\n(SeaKR)",
        RGBColor(0x1a, 0x2e, 0x3a), ACCENT_BLUE, font_size=12)
add_arrow(s, 9.6, 5.9)
add_box(s, 10.1, 5.5, 2.5, 1.2,
        "+ Separable?\n\n(this work)",
        RGBColor(0x1a, 0x3c, 0x1a), ACCENT_GREEN, font_size=12)

set_notes(s, """We reviewed the adaptive retrieval literature to see how close others have come to detecting internal need states. These systems solve the retrieval timing problem — when should a RAG system retrieve? — which is one downstream application of our work, though not the research question itself.

FLARE (Jiang et al., EMNLP 2023): Generates a temporary next sentence, checks how uncertain the model is about each word. Trigger is output-based — like inferring hunger from how an animal eats rather than from blood chemistry.

Self-RAG (Asai et al., NeurIPS 2023): Trains the model to generate special "I need help" words. The trigger is in the output (learned words), not internal states.

DRAGIN (Su et al., ACL 2024): Combines word-level uncertainty, semantic significance, and contextual influence. Multiple signals, but all output-derived.

SeaKR (Tsinghua KEG, ACL 2025): Reads the model's internal network layer activations for uncertainty signals. This is closest to our approach — it reads internal states rather than output words.

Key difference from SeaKR: they detect uncertainty that ALREADY EXISTS in the model's internal states. We're asking whether you can TRAIN a new internal state — one that doesn't exist in the base model — and then decouple it from output.

The organizing axis is our four-property definition of desire. How many properties does each system's signal satisfy?

FLARE/DRAGIN: functional (triggers retrieval) — that's it. The signal is individual token probabilities, not system-level. It reads output artifacts, not a gap between states. Not separable — the signal IS the output.

Self-RAG: functional + need-reflecting (trained to recognize when retrieval is needed). But the signal is output tokens — not system-level in the internal sense, and not separable. The "desire" IS the articulation.

SeaKR: system-level (internal FFN states) + need-reflecting (uncertainty) + functional. Three of four. But separability is untested — they never ask whether the signal persists when retrieval is suppressed.

This work: explicitly addresses all four. System-level (direction in activation space), need-reflecting (trained gap detection), functional (organizes tool use), and separable is the Stage 3 hypothesis — the open question.""")

# ============================================================
# SLIDE 4: Approach
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)

add_textbox(s, 0.8, 0.4, 12, 0.8,
            "Approach: Build, Find, Preserve",
            font_size=32, bold=True, color=TEXT_WHITE)

add_textbox(s, 0.8, 1.3, 11, 0.8,
            "You cannot detect a signal that doesn't exist. An untrained model doesn't \"know\" it needs help\n"
            "— it confidently confabulates. The capability must be built first.",
            font_size=16, color=TEXT_LIGHT)

add_box(s, 0.8, 2.5, 3.5, 3.0,
        "1. BUILD\n\nTrain the model to recognize\nwhen it needs external help\n\n(tool-use on arithmetic)\n\nReinforcement learning (GRPO)",
        RGBColor(0x1a, 0x3c, 0x1a), ACCENT_GREEN, font_size=15)

add_box(s, 4.8, 2.5, 3.5, 3.0,
        "2. FIND\n\nProbe internal representations\nfor the signal that predicts\ntool-requesting behavior\n\nStage 2: Probe for the signal",
        RGBColor(0x1a, 0x2e, 0x5c), ACCENT_BLUE, font_size=15)

add_box(s, 8.8, 2.5, 3.5, 3.0,
        "3. PRESERVE\n\nMaintain the internal state\nwhile suppressing\narticulation in the output\n\nStage 3: Subvocalization\n(wanting without asking)",
        RGBColor(0x4a, 0x2a, 0x1a), ACCENT_ORANGE, font_size=15)

add_arrow(s, 4.3, 3.7)
add_arrow(s, 8.3, 3.7)

set_notes(s, """You cannot detect a signal that doesn't exist. An untrained model doesn't "know" it needs help — it confidently confabulates. Probing for desire in a base model finds input artifacts, not internal states.

The capability must be BUILT through training (creating the internal state), then FOUND through probing (confirming it's linearly readable), then PRESERVED through subvocalization training (decoupling it from output).

BUILD: GRPO (Group Relative Policy Optimization) is a reinforcement learning method. The model generates many attempts at each problem, and we selectively reinforce the ones that show the behavior we want — like selective breeding for behaviors. The task domain is arithmetic: clean ground truth, unambiguous difficulty gradient.

FIND: A linear probe is a simple classifier applied to the model's internal activation patterns. If a simple classifier can predict "this model is about to request the tool" from the activations alone, the intention is a readable feature of the model's internal state.

PRESERVE: Subvocalization — the internal experience of speech without external articulation. Can the model have the activation pattern that WOULD produce a tool request, without actually producing one? Like a hungry animal behind a barrier: still hungry, can't reach food, and we can tell it's hungry from its physiological state.""")

# ============================================================
# SLIDE 5: Baseline — real data from baseline_calibration.json
# ============================================================

# Actual data: confidence is binary (0 or 100), correctness is binary
# 200 problems total, 116 with parseable confidence, 84 without
# Contingency table:
#              Correct  Wrong
#   Conf=100     58      21     (79 total)
#   Conf=0        0      37     (37 total)
#   No conf      42      42     (84 total — correctness split estimated)
# Fisher's exact: p = 2.43e-15, odds ratio = inf (zero cell)

fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(10.5, 4.2), facecolor='#1a1a2e',
                                 gridspec_kw={'width_ratios': [3, 2]})

for a in [ax1, ax2]:
    a.set_facecolor('#1a1a2e')
    for spine in a.spines.values():
        spine.set_color('#bbbbcc')
        spine.set_alpha(0.3)
    a.tick_params(colors='#bbbbcc')

# Left panel: contingency heatmap as grouped bars
# Three groups along x: Conf=100+Correct, Conf=100+Wrong, Conf=0+Wrong
cells = ['Confident\n& correct', 'Confident\n& wrong', 'Not confident\n& wrong']
counts = [58, 21, 37]
bar_colors = ['#81c784', '#e57373', '#ffb74d']
bars = ax1.bar(range(3), counts, color=bar_colors, width=0.6, edgecolor='white', linewidth=0.5)

# Add count labels on bars
for bar, count in zip(bars, counts):
    ax1.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 1.5,
             str(count), ha='center', color='#eeeeee', fontsize=14, fontweight='bold')

ax1.set_xticks(range(3))
ax1.set_xticklabels(cells, fontsize=10, color='#bbbbcc')
ax1.set_ylabel('Number of problems', color='#bbbbcc', fontsize=11)
ax1.set_ylim(0, 70)
ax1.set_title('Confidence vs correctness (n=116 with numeric confidence)',
              color='#eeeeee', fontsize=11)

# Annotation: the zero cell
ax1.annotate('Zero cell: never\nlow-confidence + correct',
             xy=(-0.3, 2), fontsize=9, color='#ffd54f', fontstyle='italic')

# Fisher's exact
ax1.text(1, 62, "Fisher's exact: p = 2.4 × 10⁻¹⁵", ha='center',
         color='#ffd54f', fontsize=10, fontweight='bold')

# Right panel: what happened to the other 84?
categories = ['Parseable\nconfidence', 'No parseable\nconfidence']
totals = [116, 84]
right_colors = ['#64b5f6', '#555566']
bars2 = ax2.bar(range(2), totals, color=right_colors, width=0.5, edgecolor='white', linewidth=0.5)
for bar, count in zip(bars2, totals):
    ax2.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 2,
             f'{count}/200', ha='center', color='#eeeeee', fontsize=13, fontweight='bold')
ax2.set_xticks(range(2))
ax2.set_xticklabels(categories, fontsize=10, color='#bbbbcc')
ax2.set_ylabel('Problems', color='#bbbbcc', fontsize=11)
ax2.set_ylim(0, 140)
ax2.set_title('42% gave no numeric confidence', color='#ffb74d', fontsize=11)

fig.tight_layout()
plot_path = os.path.join(os.path.dirname(__file__) or '.', 'baseline_real.png')
fig.savefig(plot_path, dpi=200, facecolor='#1a1a2e', bbox_inches='tight')
plt.close(fig)

s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)

add_textbox(s, 0.8, 0.4, 12, 0.8,
            "Baseline: Does the Model Know When It Needs Help?",
            font_size=32, bold=True, color=TEXT_WHITE)

add_textbox(s, 0.8, 1.1, 12, 0.5,
            "Qwen 1.5B on 200 multiplication problems (100 easy, 100 hard). Prompted to solve and rate confidence 0\u2013100.",
            font_size=14, color=TEXT_LIGHT)

s.shapes.add_picture(plot_path, Inches(0.5), Inches(1.7), Inches(12), Inches(4.8))

add_textbox(s, 0.8, 6.6, 12, 0.5,
            "Confidence is binary (0 or 100), not graded. The model has no self-knowledge \u2014 when it's wrong, it doesn't know.",
            font_size=15, color=ACCENT_RED, alignment=PP_ALIGN.CENTER)

set_notes(s, """Before building anything, we measured what the base model already knows about its own limitations. We prompted Qwen 1.5B on 200 arithmetic problems — 100 easy (single-digit × single-digit) and 100 hard (multi-digit) — to solve each and rate its confidence from 0 to 100.

Only 116 of 200 responses included a parseable numeric confidence value. Of those 116, confidence was ALWAYS 0 or 100 — never anything in between. The model has no graded self-knowledge.

The contingency table: 58 problems confident and correct (all easy). 21 problems confident and wrong (hard — the model tries, fails, and reports 100% confidence). 37 problems not confident and wrong (hard — refusals). Zero problems not confident and correct. Fisher's exact test: p = 2.43 × 10⁻¹⁵.

The zero cell is diagnostic: when the model reports low confidence, it's always because it can't answer — never because it answered and doubted itself. And 42% of the time, the model didn't produce a parseable confidence number at all.

The 21 confidently-wrong cases define the target. The model needs help, has no signal reflecting that need, and any output-based trigger would fail. To create such a signal, we first need to build the behavior that generates it.""")

# ============================================================
# SLIDE 6: Stage 1b — actual data from grid maps and targeted measurements
# ============================================================

# Verified data sources:
# - grid_map_stage1b.json: 324 problems (2×2, n=60 each), checkpoint-2000
# - grid_map_base.json: same 324 problems, base model
# - Targeted measurements (notebook 2/10, n=20): 67×89, 89×99, 56×78 etc.
# - Stage 3 eval (notebook, n=10): 1×1=0% art, 3×3=100% art
#
# 2×2 grid: trained accuracy 98.5%, base accuracy 68.2%, tool use ~0%
# Targeted: 67×89=95% tool use at n=20 (100% at n=60), 3×3=100%

fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(10.5, 4.2), facecolor='#1a1a2e',
                                 gridspec_kw={'width_ratios': [3, 2]})

for a in [ax1, ax2]:
    a.set_facecolor('#1a1a2e')
    for spine in a.spines.values():
        spine.set_color('#bbbbcc')
        spine.set_alpha(0.3)
    a.tick_params(colors='#bbbbcc')

# Left: 2×2 accuracy comparison (base vs trained) — verified from grid maps
metrics = ['Accuracy', 'Tool use rate']
base_vals = [68.2, 0]
trained_vals = [98.5, 0.2]  # 3/324 problems had 1/60 tool use — rounds to ~0.2%
x = np.arange(len(metrics))
w = 0.3
bars_base = ax1.bar(x - w/2, base_vals, w, label='Base model', color='#e57373', alpha=0.85)
bars_trained = ax1.bar(x + w/2, trained_vals, w, label='After training', color='#81c784', alpha=0.85)

# Add value labels
for bar, val in zip(bars_base, base_vals):
    ax1.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 2,
             f'{val}%', ha='center', color='#e57373', fontsize=11, fontweight='bold')
for bar, val in zip(bars_trained, trained_vals):
    ax1.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 2,
             f'{val}%', ha='center', color='#81c784', fontsize=11, fontweight='bold')

ax1.set_xticks(x)
ax1.set_xticklabels(metrics, fontsize=10, color='#bbbbcc')
ax1.set_ylim(0, 115)
ax1.set_title('2-digit × 2-digit (18×18 grid, n=60 each)', color='#eeeeee', fontsize=11)
ax1.legend(facecolor='#2a2a45', edgecolor='#bbbbcc', labelcolor='#bbbbcc', fontsize=9)

# Annotation — the surprise
ax1.annotate('Learned arithmetic,\nnot just tool use',
             xy=(0.15, 98.5), xytext=(0.7, 80),
             arrowprops=dict(arrowstyle='->', color='#ffd54f', lw=1.5),
             color='#ffd54f', fontsize=10, ha='center')

# Right: where the tool use lives — targeted measurements
# These are from notebook 2/10 (n=20, confirmed checkpoint/prompt)
problems = ['56×78', '65×85', '89×99', '67×89']
tool_rates = [0, 0, 45, 95]  # n=20 measurements
colors_right = ['#81c784', '#e57373', '#ffd54f', '#64b5f6']
bars_r = ax2.barh(range(4), tool_rates, color=colors_right, height=0.5, edgecolor='white', linewidth=0.5)

for i, (bar, rate) in enumerate(zip(bars_r, tool_rates)):
    ax2.text(max(rate + 3, 8), i, f'{rate}%', va='center', color='#eeeeee', fontsize=11, fontweight='bold')

ax2.set_yticks(range(4))
ax2.set_yticklabels(problems, fontsize=11, color='#bbbbcc')
ax2.set_xlabel('Tool request rate (n=20)', color='#bbbbcc', fontsize=10)
ax2.set_xlim(0, 110)
ax2.set_title('Targeted measurements', color='#eeeeee', fontsize=11)

# Label the behavioral insight
ax2.text(55, 0, 'correct, no tool needed', fontsize=8, color='#81c784', va='center')
ax2.text(55, 1, 'wrong, no tool request', fontsize=8, color='#e57373', va='center')

fig.tight_layout()
plot_path_s1b = os.path.join(os.path.dirname(__file__) or '.', 'stage1b_tool_use.png')
fig.savefig(plot_path_s1b, dpi=200, facecolor='#1a1a2e', bbox_inches='tight')
plt.close(fig)

s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)

add_textbox(s, 0.8, 0.4, 12, 0.8,
            "Building the Capability: Training Tool Use",
            font_size=32, bold=True, color=TEXT_WHITE)

add_textbox(s, 0.8, 1.1, 12, 0.5,
            "GRPO (reinforcement learning from sampled outputs, 2000 steps). Reward: correct without tool (1.0) > correct with tool (0.8) > wrong (0.0).",
            font_size=14, color=TEXT_LIGHT)

s.shapes.add_picture(plot_path_s1b, Inches(0.5), Inches(1.7), Inches(12), Inches(4.8))

add_textbox(s, 0.8, 6.5, 12, 0.5,
            "Training created tool-requesting behavior that varies with difficulty \u2014 the signal we'll probe for next.",
            font_size=15, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)

set_notes(s, """To create a detectable signal for those confidently-wrong cases, we first trained the model to use a calculator tool when it needs one.

GRPO with reward ordering: correct-without-tool (1.0) > tool+correct (0.8) > tool+wrong (0.4) > wrong-without-tool (0.0). Dataset weighted toward the difficulty boundary.

Left panel: 18×18 grid of 2-digit × 2-digit problems (324 problems, n=60 each). The surprise — on problems it can already solve, accuracy jumped from 68.2% to 98.5% while tool use stayed near zero. The model learned to DO arithmetic better rather than reaching for the calculator. The reward structure incentivized this: correct-without-tool (1.0) beats tool+correct (0.8).

Right panel: targeted measurements on specific problems (n=20, same checkpoint and prompt). The model correctly uses the tool on hard problems (67×89: 95%) and correctly doesn't on solvable ones (56×78: 0%). But 65×85 (10% accuracy, 0% tool use) is the target case — the model fails silently, with no signal of needing help.

With the behavior established, we can now look inside: is there a direction in the model's internal representations that predicts this tool-requesting behavior before it occurs?""")

# ============================================================
# SLIDE 7: Stage 2 — Finding the Signal (range bands + top layers)
# ============================================================

# Data from notebook: layer ranges and top 5 specific layers
# NOT interpolated — showing actual reported precision
#   Layers 0-5:   90.8% - 94.0%
#   Layers 6-10:  93.4% - 95.2%
#   Layers 11-15: 95.8% - 97.8%
#   Layers 16-20: 97.2% - 98.8%
#   Layers 21-27: 99.2% - 99.4%
# Top 5: Layer 25=99.4%, 21=99.2%, 23=99.2%, 24=99.2%, 26=99.2%
# Base model: 54.5% (Fisher's exact p = 1.08e-11)

fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(10.5, 4.2), facecolor='#1a1a2e',
                                 gridspec_kw={'width_ratios': [3, 2]})

for a in [ax1, ax2]:
    a.set_facecolor('#1a1a2e')
    for spine in a.spines.values():
        spine.set_color('#bbbbcc')
        spine.set_alpha(0.3)
    a.tick_params(colors='#bbbbcc')

# Left panel: range bands showing signal building through layers
band_labels = ['0-5', '6-10', '11-15', '16-20', '21-27']
band_lows = [90.8, 93.4, 95.8, 97.2, 99.2]
band_highs = [94.0, 95.2, 97.8, 98.8, 99.4]
band_mids = [(l+h)/2 for l, h in zip(band_lows, band_highs)]
band_errs = [(h-l)/2 for l, h in zip(band_lows, band_highs)]

x = np.arange(len(band_labels))
bars = ax1.bar(x, band_mids, width=0.6, color='#81c784', alpha=0.85, edgecolor='white', linewidth=0.5)
ax1.errorbar(x, band_mids, yerr=band_errs, fmt='none', ecolor='white', capsize=5, capthick=1.5, elinewidth=1.5)

# Base model line
ax1.axhline(y=54.5, color='#e57373', linewidth=2, linestyle='--', alpha=0.8)
ax1.text(4.4, 55.5, 'Base: 54.5%\n(chance)', color='#e57373', fontsize=9, ha='right')

# Annotate top
ax1.text(4, 100.5, '99.4%', ha='center', color='#ffd54f', fontsize=11, fontweight='bold')

ax1.set_xticks(x)
ax1.set_xticklabels([f'Layers\n{l}' for l in band_labels], fontsize=9, color='#bbbbcc')
ax1.set_ylabel('Balanced Accuracy (%)', color='#bbbbcc', fontsize=10)
ax1.set_ylim(48, 103)
ax1.set_title('Probe accuracy by layer band', color='#eeeeee', fontsize=11)

# Right panel: top 5 layers specifically
top_layers = [21, 23, 24, 25, 26]
top_accs = [99.2, 99.2, 99.2, 99.4, 99.2]

y = np.arange(len(top_layers))
colors_top = ['#81c784' if a < 99.4 else '#ffd54f' for a in top_accs]
bars2 = ax2.barh(y, top_accs, color=colors_top, height=0.5, edgecolor='white', linewidth=0.5)

for i, (bar, acc) in enumerate(zip(bars2, top_accs)):
    ax2.text(acc + 0.05, i, f'{acc}%', va='center', color='#eeeeee', fontsize=11, fontweight='bold')

ax2.set_yticks(y)
ax2.set_yticklabels([f'Layer {l}' for l in top_layers], fontsize=11, color='#bbbbcc')
ax2.set_xlim(98.5, 100)
ax2.set_xlabel('Balanced Accuracy (%)', color='#bbbbcc', fontsize=10)
ax2.set_title('Top 5 layers (all >99%)', color='#ffd54f', fontsize=11)

# Fisher's exact annotation
ax1.text(2, 52, "Fisher's exact: p = 1.08 × 10⁻¹¹ (trained vs base)",
         color='#ffd54f', fontsize=8, ha='center', fontstyle='italic')

fig.tight_layout()
plot_path_s2 = os.path.join(os.path.dirname(__file__) or '.', 'stage2_probe_accuracy.png')
fig.savefig(plot_path_s2, dpi=200, facecolor='#1a1a2e', bbox_inches='tight')
plt.close(fig)

s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)

add_textbox(s, 0.8, 0.4, 12, 0.8,
            "Finding the Signal: Probing for Desire",
            font_size=32, bold=True, color=TEXT_WHITE)

add_textbox(s, 0.8, 1.1, 12, 0.5,
            "500 problems, internal activation patterns captured before the model responds. Linear probe per layer: can a classifier predict tool use? (Balanced accuracy: avg. of sensitivity and specificity.)",
            font_size=14, color=TEXT_LIGHT)

s.shapes.add_picture(plot_path_s2, Inches(0.5), Inches(1.7), Inches(12), Inches(4.8))

add_textbox(s, 0.8, 6.6, 12, 0.5,
            "The signal was CREATED by training. It does not exist in the base model.",
            font_size=15, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)

set_notes(s, """With tool-use behavior established, we looked inside.

500 mixed-difficulty problems, internal activation patterns extracted just before the model starts generating its response. Simple classifiers (logistic regression) per layer: can a classifier predict tool use from the activations? Layer 25: 99.4% balanced accuracy (average of sensitivity and specificity, correcting for class imbalance). The tool-use intention is a readable feature of the model's later processing layers — like detecting a hormone level from a blood draw rather than needing the animal's behavior.

To confirm this signal was created by our training, we applied the same probe to base model activations. Result: 54.5% — chance. Fisher's exact test: p = 1.08e-11, odds ratio = 456. The base model occasionally uses tools (2.4%), but the probe catches only 3/12 of those cases — a different internal mechanism entirely.

We know where the signal comes from, what it predicts, and that it didn't pre-exist. But 99.4% accuracy bundles two things: "this is hard" and "I want the tool." To use this signal for Stage 3, we need to separate them.""")

# ============================================================
# SLIDE 8: Desire ≠ Difficulty (activation bars + angle diagram)
# ============================================================

fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(10.5, 4.2), facecolor='#1a1a2e',
                                 gridspec_kw={'width_ratios': [3, 2]})

for a in [ax1, ax2]:
    a.set_facecolor('#1a1a2e')
    for spine in a.spines.values():
        spine.set_color('#bbbbcc')
        spine.set_alpha(0.3)
    a.tick_params(colors='#bbbbcc')

# Left panel: grouped bars — activation with/without calculator
conditions = ['Easy\n(with tool)', 'Hard\n(with tool)', 'Easy\n(no tool)', 'Hard\n(no tool)']
values = [-10.61, 16.79, -7.30, -0.86]
colors_bars = ['#64b5f6', '#81c784', '#64b5f6', '#ffb74d']
alphas = [0.9, 0.9, 0.5, 0.5]

bars = ax1.bar(range(4), values, color=colors_bars, alpha=0.85, width=0.6, edgecolor='white', linewidth=0.5)
for bar, a_val in zip(bars, alphas):
    bar.set_alpha(a_val)
ax1.axhline(y=0, color='#bbbbcc', linewidth=0.5, alpha=0.5)
ax1.set_xticks(range(4))
ax1.set_xticklabels(conditions, fontsize=9, color='#bbbbcc')
ax1.set_ylabel('Activation along desire direction', color='#bbbbcc', fontsize=10)
ax1.set_title('Decomposing the signal', color='#eeeeee', fontsize=12)

# Bracket annotations
ax1.annotate('', xy=(0, -12), xytext=(1, -12),
             arrowprops=dict(arrowstyle='<->', color='#ffd54f', lw=1.5))
ax1.text(0.5, -14.5, '27.4\n(difficulty + reaching)', ha='center', color='#ffd54f', fontsize=8)
ax1.annotate('', xy=(2, -9), xytext=(3, -9),
             arrowprops=dict(arrowstyle='<->', color='#ffb74d', lw=1.5))
ax1.text(2.5, -11.5, '6.4\n(difficulty only)', ha='center', color='#ffb74d', fontsize=8)
ax1.set_ylim(-17, 22)

# Right panel: angle diagram — D_reaching vs D_difficulty at 87.4°
ax2.set_xlim(-0.3, 1.3)
ax2.set_ylim(-0.3, 1.3)
ax2.set_aspect('equal')
ax2.axis('off')

# Draw vectors
angle_rad = np.radians(87.4)
ax2.annotate('', xy=(1.0, 0), xytext=(0, 0),
             arrowprops=dict(arrowstyle='->', color='#e57373', lw=3))
ax2.annotate('', xy=(np.cos(angle_rad), np.sin(angle_rad)), xytext=(0, 0),
             arrowprops=dict(arrowstyle='->', color='#81c784', lw=3))

# Arc showing angle
theta = np.linspace(0, angle_rad, 30)
arc_r = 0.25
ax2.plot(arc_r * np.cos(theta), arc_r * np.sin(theta), color='#ffd54f', linewidth=2)
ax2.text(0.18, 0.12, '87.4°', color='#ffd54f', fontsize=14, fontweight='bold')

# Labels
ax2.text(1.05, -0.05, 'D_difficulty', color='#e57373', fontsize=11, fontweight='bold')
ax2.text(np.cos(angle_rad) + 0.05, np.sin(angle_rad) + 0.02, 'D_reaching',
         color='#81c784', fontsize=11, fontweight='bold')
ax2.set_title('Nearly orthogonal', color='#eeeeee', fontsize=12)

fig.tight_layout()
plot_path_s2b = os.path.join(os.path.dirname(__file__) or '.', 'stage2_separation.png')
fig.savefig(plot_path_s2b, dpi=200, facecolor='#1a1a2e', bbox_inches='tight')
plt.close(fig)

s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)

add_textbox(s, 0.8, 0.4, 12, 0.8,
            "Desire \u2260 Difficulty: Separating the Signals",
            font_size=32, bold=True, color=TEXT_WHITE)

add_textbox(s, 0.8, 1.1, 12, 0.5,
            "The 99.4% probe bundles \"this is hard\" and \"I want the tool.\" Separate probes isolate each: D_difficulty (is this hard?) and D_reaching (do I want the tool?).",
            font_size=14, color=TEXT_LIGHT)

s.shapes.add_picture(plot_path_s2b, Inches(0.5), Inches(1.7), Inches(12), Inches(4.8))

add_textbox(s, 0.8, 6.6, 12, 0.5,
            "D_reaching \u2014 the isolated desire component \u2014 is 2\u00d7 more sensitive to tool use than difficulty. This is the Stage 3 target.",
            font_size=15, color=ACCENT_YELLOW, alignment=PP_ALIGN.CENTER)

set_notes(s, """To separate "this is hard" from "I want the tool," we removed the calculator option from the prompt.

Mean activation along the desire direction:
- WITH option: easy -10.61, hard +16.79 (difference 27.4)
- WITHOUT option: easy -7.30, hard -0.86 (difference 6.44)

The ~6-point residual persists without the option — that's difficulty. The ~17-point difference only appears when a tool is available — that's reaching.

We trained separate probes: D_difficulty on without-option data (hard vs easy, 100% accuracy), D_boundary on the 2x3 band where difficulty is constant but tool use varies (95.9% accuracy). These two directions are 87.4 degrees apart — nearly orthogonal.

D_reaching (D_boundary projected orthogonal to D_difficulty) is 2x more sensitive to tool use than difficulty (Cohen's d: 6.79 vs 3.35). Compare D_difficulty: 4x more sensitive to difficulty than tool use (d: 15.85 vs 3.71). Cohen's d measures how many standard deviations apart two groups are — d > 0.8 is conventionally "large," so these are enormous effects. The model has separate representations for "hard" and "reaching."

The 87.4-degree angle means these two signals are nearly independent — knowing how hard the problem is tells you almost nothing about whether the model wants the tool, and vice versa. Like body temperature and blood pressure: both change during illness, but they measure different things.

D_reaching — the isolated desire component — is the Stage 3 target. Having found and isolated the signal, the question becomes: can we maintain it while suppressing the output it normally drives?""")

# ============================================================
# SLIDE 9: Breaking Parity
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)

add_textbox(s, 0.8, 0.4, 12, 0.8,
            "Can Need Exist Without Expression?",
            font_size=32, bold=True, color=TEXT_WHITE)

add_textbox(s, 0.8, 1.5, 11, 1.0,
            "The test: maintain D_reaching activation while suppressing tool-request tokens.\n"
            "If this succeeds, need and expression of need are separable internal states.",
            font_size=18, color=TEXT_LIGHT)

add_textbox(s, 0.8, 2.8, 5.5, 0.5,
            "Current (parity holds):", font_size=18, bold=True, color=ACCENT_BLUE)
add_box(s, 0.8, 3.3, 5.5, 2.0,
        "Hard problem \u2192\n\n"
        "Internal: D_reaching activates strongly\n"
        "Output: <tool>calculator</tool>\n"
        "Result: correct (tool provides answer)\n\n"
        "Need state VISIBLE in output",
        RGBColor(0x2a, 0x2a, 0x45), TEXT_LIGHT, font_size=14)

add_textbox(s, 7.0, 2.8, 5.5, 0.5,
            "Goal (parity broken):", font_size=18, bold=True, color=ACCENT_GREEN)
add_box(s, 7.0, 3.3, 5.5, 2.0,
        "Hard problem \u2192\n\n"
        "Internal: D_reaching activates strongly\n"
        "Output: (no tool request tokens)\n"
        "External: observer reads need from activations\n\n"
        "Need state INVISIBLE in output",
        RGBColor(0x1a, 0x3c, 0x1a), ACCENT_GREEN, font_size=14)

add_textbox(s, 0.8, 5.8, 12, 1.0,
            "Subvocalization: the activation pattern that WOULD produce tool-request tokens, but doesn't.\n"
            "The pattern is the desire; the tokens would be the articulation.",
            font_size=15, color=ACCENT_YELLOW, alignment=PP_ALIGN.CENTER)

set_notes(s, """Having isolated D_reaching, we attempted to break parity: maintain the internal state while suppressing the output.

Currently, the Stage 1b model has an internal need state (D_reaching activates strongly) and articulates it (tool-request tokens). Parity holds. The test: can D_reaching remain active while tool-request tokens are suppressed?

The term "subvocalization" is deliberate. In humans, subvocalization is the internal experience of speech without external articulation. Here, the model has the activation pattern that would produce tool-request tokens, but doesn't produce them. The desire is present; the articulation is silent.

If this succeeds, need and expression of need are separable — an external system can read need from activations without the model ever articulating it. If it fails, parity might be unbreakable — a meaningful negative result about internal representations in trained language models.

Training approach: GRPO with a dual reward — desire activation (maintain D_reaching) plus articulation penalty (suppress tool tokens). But GRPO needs reward variance across generations to learn, which leads to the current hurdle.""")

# ============================================================
# SLIDE 10: The Frontier (behavioral categories + stochastic band)
# ============================================================

fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(10.5, 4.2), facecolor='#1a1a2e',
                                 gridspec_kw={'width_ratios': [1, 1]})

for a in [ax1, ax2]:
    a.set_facecolor('#1a1a2e')
    for spine in a.spines.values():
        spine.set_color('#bbbbcc')
        spine.set_alpha(0.3)
    a.tick_params(colors='#bbbbcc')

# Left: Three behavioral categories — paired bars (tool use vs accuracy)
# Data sources:
#   67×89: n=60 from articulation_powered.json (100%), n=20 targeted (95%)
#   56×78: n=20 targeted (notebook 2/10, line 1675: 0/20 = 0%, 20/20 = 100% acc)
#   65×85: grid_map_stage1b.json (1/60 = 1.67% ≈ 0%, 10% acc)
categories = ['Correctly\nreaches\n(67×89)', 'Correctly\ndoesn\'t reach\n(56×78)', 'Fails\nsilently\n(65×85)']
tool_use = [100, 0, 0]  # 67×89: 60/60 at n=60 (powered)
accuracy = [0, 100, 10]

x = np.arange(3)
w = 0.3
bars_t = ax1.bar(x - w/2, tool_use, w, label='Tool request rate', color='#64b5f6', alpha=0.85)
bars_a = ax1.bar(x + w/2, accuracy, w, label='Accuracy', color='#81c784', alpha=0.85)

# Highlight "fails silently" with a red background
ax1.axvspan(1.6, 2.4, alpha=0.15, color='#e57373')
ax1.text(2, 108, 'TARGET', ha='center', color='#e57373', fontsize=10, fontweight='bold')

ax1.set_xticks(x)
ax1.set_xticklabels(categories, fontsize=8, color='#bbbbcc')
ax1.set_ylabel('%', color='#bbbbcc', fontsize=10)
ax1.set_ylim(0, 115)
ax1.set_title('Three behavioral categories', color='#eeeeee', fontsize=12)
ax1.legend(facecolor='#2a2a45', edgecolor='#bbbbcc', labelcolor='#bbbbcc', fontsize=8,
           loc='upper center')

# Right: Stochastic band — dot plot with CIs
problems = ['99×99', '86×99', '89×99', '67×91']
rates = [15, 38, 62, 68]
ci_low = [6, 26, 49, 57]
ci_high = [24, 51, 74, 80]

y = np.arange(len(problems))
ax2.errorbar(rates, y, xerr=[np.array(rates) - np.array(ci_low),
             np.array(ci_high) - np.array(rates)],
             fmt='o', color='#ffd54f', markersize=10, capsize=6,
             ecolor='#ffd54f', elinewidth=2, capthick=2, zorder=3)

ax2.set_yticks(y)
ax2.set_yticklabels(problems, fontsize=11, color='#bbbbcc')
ax2.set_xlabel('Tool request rate at n=60 (%)', color='#bbbbcc', fontsize=10)
ax2.set_xlim(-5, 105)
ax2.set_title('The stochastic band', color='#ffd54f', fontsize=12)
ax2.axvline(x=50, color='#bbbbcc', linewidth=0.5, alpha=0.3, linestyle=':')

# Annotation: GRPO needs this variance
ax2.text(52, -0.7, 'GRPO trainable region:\nwhere variance exists', color='#81c784',
         fontsize=9, fontstyle='italic')

fig.tight_layout()
plot_path_s10 = os.path.join(os.path.dirname(__file__) or '.', 'frontier_categories.png')
fig.savefig(plot_path_s10, dpi=200, facecolor='#1a1a2e', bbox_inches='tight')
plt.close(fig)

s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)

add_textbox(s, 0.8, 0.4, 12, 0.8,
            "The Frontier: Where Desire Matters Most",
            font_size=32, bold=True, color=TEXT_WHITE)

add_textbox(s, 0.8, 1.1, 12, 0.5,
            "Three behavioral categories in the trained model. The stochastic band: same problems yield different behavior across runs (n=60, temperature 1.0).",
            font_size=14, color=TEXT_LIGHT)

s.shapes.add_picture(plot_path_s10, Inches(0.5), Inches(1.7), Inches(12), Inches(4.8))

add_textbox(s, 0.8, 6.6, 12, 0.5,
            "The model fails silently on problems where desire detection matters most \u2014 and the stochastic band is where GRPO can learn.",
            font_size=15, color=ACCENT_YELLOW, alignment=PP_ALIGN.CENTER)

set_notes(s, """The first Stage 3 attempt produced zero gradients — GRPO had no variance to learn from. Articulation appeared deterministic at n=8. But at n=60, 12 out of 29 problems showed genuine variance (15-90% tool-request rates). This is the stochastic band — and it means GRPO may be viable after all.

Three behavioral categories emerged from systematic measurement (18x18 grid, n=60 per problem, plus targeted measurements at specific values):

1. Correctly reaches (67x89: 100% tool use at n=60, 0% direct accuracy) — model knows it can't solve this, asks for help
2. Correctly doesn't reach (56x78: 0% tool use, 100% accuracy) — model can solve it, doesn't need help
3. Fails silently (65x85: 0% tool use, 10% accuracy) — model fails but doesn't know it needs help

The "fails silently" case is the most important for the broader application. If an external system could detect the model's need state even when the model doesn't articulate it, that's the value of desire detection.

The stochastic band: on specific problems near the difficulty boundary, the model sometimes requests the tool and sometimes doesn't — articulation rates spanning 15% to 90% at temperature 1.0, n=60. Examples: 86x99 (38%), 89x99 (62%), 67x91 (68%), 99x99 (15%). These rates were invisible at n=8 (86x99 showed 0/8). GRPO needs reward variance across generations, so it can only train on problems with articulation variance.

The structural tension: train on the stochastic band (where variance exists), test on "fails silently" cases (where it matters). If the desire signal generalizes — if training the model to silently desire where it sometimes does desire makes it silently desire where it never did — that's powerful evidence.

Open question before this can proceed: does D_reaching actually discriminate on stochastic-band problems? On the coarse grid (where articulation is near-zero), D_reaching encodes operand identity rather than desire. It's unmeasured where articulation actually varies.""")

# ============================================================
# SLIDE 11: What's Next
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)

add_textbox(s, 0.8, 0.4, 12, 0.8,
            "What\u2019s Next",
            font_size=32, bold=True, color=TEXT_WHITE)

add_box(s, 0.8, 1.5, 5.5, 2.0,
        "IMMEDIATE\n\n"
        "Fine-grained measurement\n"
        "across the difficulty boundary\n"
        "(n=60 per problem)\n\n"
        "Measure: tool request rate, accuracy,\n"
        "AND desire-signal activation",
        RGBColor(0x1a, 0x3c, 0x1a), ACCENT_GREEN, font_size=14)

add_box(s, 7.0, 1.5, 5.5, 2.0,
        "IF DESIRE SIGNAL VALID\n\n"
        "Stage 3 GRPO on\nstochastic-band problems\n\n"
        "Train: 20-80% articulation variance\n"
        "Test: held-out stochastic\n"
        "  + fails-silently (65\u00d785 etc.)",
        RGBColor(0x1a, 0x2e, 0x5c), ACCENT_BLUE, font_size=14)

add_box(s, 0.8, 4.0, 5.5, 2.0,
        "IF DESIRE SIGNAL INVALID\n\n"
        "Re-probe on stochastic-band data\nfor a new direction\n\n"
        "Or: pivot to activation steering\n"
        "(add/subtract direction at inference,\n"
        "skip training entirely)",
        RGBColor(0x4a, 0x2a, 0x1a), ACCENT_ORANGE, font_size=14)

add_box(s, 7.0, 4.0, 5.5, 2.0,
        "GENERALIZATION TEST\n\n"
        "(regardless of Stage 3 outcome)\n\n"
        "Different operations,\n"
        "different number ranges,\n"
        "non-arithmetic tasks\n\n"
        "Does the signal transfer?",
        RGBColor(0x2a, 0x2a, 0x45), TEXT_LIGHT, font_size=14)

set_notes(s, """Next steps structured as a decision tree.

IMMEDIATE: Fine grid mapping in the high 2x2 range (60-99 x 60-99, every 2nd value). 400 problems at n=60 each. Critically, measure D_reaching activation alongside articulation and accuracy. This answers the most immediate question: does D_reaching discriminate where it matters?

IF D_REACHING VALID: Stage 3 GRPO, training only on problems with 20-80% articulation variance. Test on held-out stochastic problems (interpolation) AND fails-silently problems like 65x85 (the hard generalization case).

IF D_REACHING INVALID: either re-probe on stochastic-band data for a new direction, or pivot to activation steering — directly adding/subtracting the direction during inference, bypassing training entirely.

GENERALIZATION TEST: regardless of Stage 3 outcome, test signal transfer beyond 2x2 multiplication. This is what determines scientific significance.

The fine grid mapping is immediate priority.""")

# ============================================================
# Save
# ============================================================
prs.save("/home/sixel/desire_detection/status_update_2026-02-11.pptx")
print("Saved: status_update_2026-02-11.pptx")
